from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

def add_slide(title, content_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
    return slide

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Optimasi System Monitoring Project Menggunakan AI\nuntuk Memudahkan Evaluasi dan Meningkatkan Kinerja"
subtitle.text = "Inovasi: AIMON (AI Monitoring for Optimization and Navigation)"

# 01 Latar Belakang
add_slide("01 - Latar Belakang", [
    "Dalam pelaksanaan project, proses monitoring dan evaluasi sering kali membutuhkan waktu, ketelitian, serta analisis data yang cukup besar.",
    "Proses manual berpotensi menimbulkan keterlambatan identifikasi deviasi dan menyulitkan rekapitulasi data.",
    "Kurang optimalnya pengambilan keputusan pimpinan karena kondisi aktual project tidak dapat dipantau secara langsung.",
    "Tujuan: Mengoptimalkan pemantauan berbasis AI, mengurangi monitoring manual, dan mempercepat identifikasi hambatan operasional."
])

# 02 Masalah
add_slide("02 - Masalah", [
    "Apa yang terjadi: Proses pelaporan progres proyek (di Area OMTI, Non-OMTI, Security Printing, dll) berjalan secara konvensional dan terpecah (silo).",
    "Baseline permasalahan: Dokumen (evidence) dan laporan progres tersebar di email/aplikasi chat personal, sehingga sulit dilacak dan diaudit.",
    "Dampak: Pemborosan man-hours untuk mengumpulkan laporan, serta rantai persetujuan (approval) yang sangat lambat dari level pelaksana hingga level eksekutif."
])

# 03 Analisis Penyebab
add_slide("03 - Analisis Penyebab", [
    "Akar Masalah (Root Cause): Tidak adanya sebuah ekosistem digital terpadu yang memfasilitasi 3 pilar utama secara bersamaan: Manajemen Proyek, Komunikasi, dan Manajemen Bukti Dokumen.",
    "Dampak turunan 1: Laporan lambat direkap karena atasan harus memvalidasi data dan mencari dokumen di tempat yang berbeda-beda.",
    "Dampak turunan 2: Jika terjadi revisi, komunikasi tercecer karena tidak ada fitur diskusi/chat yang menempel langsung pada konteks proyek tersebut."
])

# 04 Solusi (Iterasi 1)
add_slide("04 - Solusi: Iterasi 1 (Fondasi Dashboard & Approval)", [
    "Langkah awal pembuatan AIMON difokuskan pada perancangan arsitektur antarmuka (UI/UX) dan sistem hierarki pengguna.",
    "Sistem Role-based: Hak akses dibedakan untuk 5 tingkat (Researcher -> Kepala Unit -> Kepala Seksi -> Kepala Departemen -> Kepala Divisi).",
    "Approval Queue: Membangun alur persetujuan berjenjang secara digital. Bawahan mengirim update, dan atasan dapat melakukan verifikasi (Approve/Reject) dalam 1 klik.",
    "Fitur Draft: Penambahan fitur 'Save as Draft' agar Researcher dapat menyimpan progres sementara sebelum disubmit."
])

# 04 Solusi (Iterasi 2)
add_slide("04 - Solusi: Iterasi 2 (Integrasi Database & Cloud)", [
    "Setelah antarmuka selesai, pengembangan dilanjutkan ke arsitektur backend untuk menjamin keamanan dan persistensi data.",
    "Cloud Database: Mengimplementasikan PostgreSQL di Cloud (Supabase) dikombinasikan dengan Prisma ORM untuk manajemen data relasional.",
    "Kredensial Aman: Mendaftarkan seluruh akun (9 User utama dari berbagai level) secara terpusat agar sistem bisa diakses dari berbagai perangkat secara konsisten tanpa kehilangan data."
])

# 04 Solusi (Iterasi 3)
add_slide("04 - Solusi: Iterasi 3 (File Management Terpusat)", [
    "Kendala pada sistem lama adalah hilangnya jejak dokumen pendukung (evidence).",
    "Pengembangan Backend: Menggunakan Node.js dan library Multer untuk menciptakan sistem 'File Upload' yang robust.",
    "Single Source of Truth: Kini, setiap update progres di AIMON diwajibkan menyertakan file asli (PDF, gambar, dokumen).",
    "Aksesibilitas: Atasan yang me-review dapat langsung mengklik 'Lihat' atau 'Unduh' file tersebut tanpa harus keluar dari dashboard."
])

# 04 Solusi (Iterasi 4)
add_slide("04 - Solusi: Iterasi 4 (Sistem Komunikasi Real-time)", [
    "Untuk mengatasi masalah lambatnya komunikasi revisi, dibangun fitur In-App Chat ala WhatsApp di dalam AIMON.",
    "Real-time Notification: Munculnya angka notifikasi (badge merah) dan highlight teks berwarna biru ketika ada pesan baru yang belum dibaca.",
    "Smart Sorting: Kontak yang baru saja mengirim pesan otomatis naik ke urutan paling atas.",
    "UI/UX Enterprise: Implementasi pop-up modal elegan untuk penolakan (reject), fitur auto-initial untuk foto profil (misal: 'TA' untuk Titi Andayani), dan status online."
])

# 05 Implementasi
add_slide("05 - Implementasi", [
    "Lingkup Uji (MVP): 9 Pengguna utama yang mewakili eksekusi proyek di berbagai area.",
    "Kendala Implementasi: Terjadi gap sinkronisasi data ketika beberapa atasan dan bawahan menggunakan sistem secara bersamaan di browser berbeda.",
    "Cara Penanggulangan: Mengembangkan arsitektur Web Event Listener berbasis Local Storage dan API backend yang sinkron secara instan antar tab/komputer tanpa perlu refresh halaman."
])

# 06 Dampak
add_slide("06 - Dampak (Impact)", [
    "Cost Avoidance: Menghilangkan puluhan jam kerja rekapitulasi manual di akhir bulan.",
    "Transparansi & Kepatuhan: Setiap tindakan (Approve, Reject, Update) memiliki rekam jejak digital yang tidak bisa dimanipulasi, lengkap dengan dokumen aslinya.",
    "Kepuasan Pengguna: Kurva pembelajaran (learning curve) sangat rendah karena UI/UX AIMON meniru aplikasi yang sudah familiar (misal: antarmuka chat mirip WhatsApp)."
])

# 07 Sustainability
add_slide("07 - Sustainability", [
    "Infrastruktur Terpusat: Seluruh sistem (Frontend dan Backend) akan di-deploy ke Server/Komputer Kantor yang menyala 24 jam penuh.",
    "Aksesibilitas 24/7: Dapat diakses oleh staf manapun dalam jaringan menggunakan IP address server, memastikan kontinuitas pemantauan.",
    "Standardisasi SOP: Sistem ini menciptakan SOP baru secara teknis (by-design), di mana sistem akan menolak update jika bukti tidak lengkap."
])

# 08 Lesson Learned
add_slide("08 - Lesson Learned", [
    "Tantangan Utama: Mengembangkan aplikasi yang harus menjembatani kebutuhan Top Management (data agregat) dan level staf operasional (data detail dan alat komunikasi).",
    "Pembelajaran Penting: Proses pengembangan yang iteratif (berubah-ubah sesuai masukan) membuktikan bahwa 'User Experience' adalah segalanya.",
    "Fitur-fitur detail (seperti badge chat, pop-up tolak elegan, inisial profil) adalah kunci yang membuat user mau meninggalkan cara manual.",
    "Kesimpulan: AIMON berhasil mentransformasi tata kelola proyek menjadi budaya kerja yang modern, transparan, dan berbasis data real-time."
])

ppt_path = "D:\\2026\\IAKA\\AIMON (AI Monitoring for Optimization and Navigation)\\Materi_Presentasi_IAKA_AIMON_Lengkap.pptx"
prs.save(ppt_path)
print(f"Berhasil membuat PPT di: {ppt_path}")
