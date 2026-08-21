from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

def add_slide(title, content_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
    return slide

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Optimasi System Monitoring Project Menggunakan AI\nuntuk Memudahkan Evaluasi dan Meningkatkan Kinerja"
subtitle.text = "Inovasi: AIMON (AI Monitoring for Optimization and Navigation)"

# Slide 2: Latar Belakang (Part 1)
add_slide("01 - Latar Belakang", [
    "Dalam pelaksanaan project, proses monitoring dan evaluasi sering kali membutuhkan waktu, ketelitian, serta analisis data yang cukup besar.",
    "Kondisi saat ini: Proses manual berpotensi menimbulkan keterlambatan dalam identifikasi deviasi.",
    "Kesulitan dalam merangkum data dari berbagai sumber komunikasi yang tidak terpusat.",
    "Akibatnya: Kurang optimalnya pengambilan keputusan berdasarkan kondisi aktual project.",
    "Tujuan Inovasi: Mengurangi ketergantungan pada proses manual dan mempermudah evaluasi."
])

# Slide 3: Masalah (Part 2)
add_slide("02 - Masalah", [
    "Apa yang terjadi: Laporan progres (update nilai) proyek dan dokumen evidence dilakukan secara terpisah (misal via WA/Email), melibatkan banyak pihak mulai dari Researcher hingga Kepala Divisi.",
    "Baseline permasalahan: Hilangnya jejak history persetujuan (approval) dan catatan revisi antar level manajemen.",
    "Waktu siklus pelaporan: Membutuhkan waktu berhari-hari hanya untuk mengumpulkan data dan dokumen ke dalam satu laporan terpadu.",
    "Dampak: Pemantauan tidak real-time, menyulitkan pimpinan dalam mengambil keputusan strategis tepat waktu."
])

# Slide 4: Analisis Penyebab (Part 3)
add_slide("03 - Analisis Penyebab", [
    "Mengapa pelaporan lambat? Karena tidak ada satu sistem (Single Source of Truth) yang mengintegrasikan dokumen, persetujuan, dan obrolan.",
    "Mengapa proses persetujuan (Approval) sering tersendat? Tidak ada sistem notifikasi langsung (real-time) ketika dokumen membutuhkan review.",
    "Akar Masalah (Root Cause): Belum adanya platform digital monitoring yang spesifik mengakomodasi alur kerja (workflow) pelaporan proyek dari bawah ke atas secara real-time dan terstruktur.",
    "Bukti pendukung: Komunikasi mengenai revisi dokumen proyek seringkali tertumpuk di aplikasi chat eksternal (silo)."
])

# Slide 5: Solusi (Part 4)
add_slide("04 - Solusi: AIMON Dashboard", [
    "Konsep Solusi: Membangun AIMON (AI Monitoring for Optimization and Navigation) sebagai dashboard terpusat berbasis web.",
    "Mekanisme Kerja:",
    "- Sistem Role-based (Researcher, Kepala Unit, Kepala Seksi, Kadep, Kadiv) dengan hak akses berjenjang.",
    "- Sistem Approval Queue: Dokumen & update persentase proyek diverifikasi secara terstruktur.",
    "- In-App Real-time Chat: Diskusi spesifik langsung di dalam sistem tanpa perlu platform eksternal.",
    "- Document Management: Semua lampiran evidence tersimpan rapi pada masing-masing tugas.",
    "Alur Before: Manual -> Chat eksternal -> Rekap Manual",
    "Alur After: Update di AIMON -> Notifikasi -> Klik Approve -> Selesai"
])

# Slide 6: Implementasi (Part 5)
add_slide("05 - Implementasi", [
    "Lingkup Uji Coba (MVP): Diterapkan untuk monitoring program di berbagai area (OMTI, Security Printing, dsb).",
    "Pengguna: Melibatkan seluruh hierarki dari Researcher (pembuat draft update) hingga level manajemen (pemberi approval).",
    "Kendala yang dihadapi: Sinkronisasi data dan chat secara real-time antar perangkat agar terasa seperti aplikasi native.",
    "Cara penanggulangan: Mengimplementasikan sistem sinkronisasi otomatis dan notifikasi visual (badge/indikator warna merah) untuk update chat dan dokumen terbaru.",
    "Data Before vs After: Waktu pengajuan update progress hingga disetujui berkurang secara drastis."
])

# Slide 7: Dampak (Part 6)
add_slide("06 - Dampak (Impact)", [
    "Cost Avoidance & Efisiensi Waktu: Mengurangi secara drastis jam kerja (man-hours) yang terbuang untuk rekapitulasi data progres manual.",
    "Mutu & Kepatuhan: Setiap perubahan persentase progress terekam dalam histori beserta dokumen buktinya (evidence-based).",
    "Kenyamanan Pengguna: Antarmuka modern yang adaptif (UI/UX layaknya aplikasi komersial) dengan fitur hapus foto profil, notifikasi chat, dan sistem draft otomatis.",
    "Potensi Replikasi: Arsitektur aplikasi web AIMON ini dirancang sangat dinamis sehingga mudah diadaptasi untuk monitoring proyek di Divisi atau Departemen lain."
])

# Slide 8: Sustainability (Part 7)
add_slide("07 - Sustainability", [
    "Sistem Terpusat: Semua kode sistem (Frontend & Backend Database Supabase) ditempatkan dan dijalankan pada Server/Komputer Kantor pusat yang menyala 24 jam.",
    "Data Permanen: Segala update akan selalu tersimpan dan dapat diakses dari perangkat manapun oleh anggota tim yang memiliki hak akses.",
    "Transfer Knowledge: Fitur yang sangat intuitif menekan kurva pembelajaran (learning curve). Standarisasi input data terjamin karena sistem menolak update jika bukti tidak lengkap."
])

# Slide 9: Lesson Learned (Part 8)
add_slide("08 - Lesson Learned", [
    "Tantangan Utama: Merancang User Experience (UX) untuk pengguna dari berbagai level hierarki (mulai dari eksekutor hingga Top Management) agar sama-sama mudah digunakan.",
    "Pembelajaran: Fitur sekecil apapun, seperti pop-up catatan penolakan (Reject Modal) yang terlihat elegan dan fitur notifikasi Chat (Real-time badge), ternyata sangat memengaruhi kemauan pengguna untuk beralih dari cara manual.",
    "Kesimpulan: Optimalisasi menggunakan AIMON bukan hanya sekadar digitalisasi, melainkan pembentukan budaya kerja baru yang lebih transparan, terukur, dan berbasis data (data-driven)."
])

ppt_path = "D:\\2026\\IAKA\\AIMON (AI Monitoring for Optimization and Navigation)\\Materi_Presentasi_IAKA_AIMON.pptx"
prs.save(ppt_path)
print(f"Berhasil membuat PPT di: {ppt_path}")
